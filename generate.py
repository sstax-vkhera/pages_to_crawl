import os
import random
import string
from io import BytesIO

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import xlsxwriter
from fpdf import FPDF
from PyPDF2 import PdfWriter
import msoffcrypto


def random_text(length=100):
    return ''.join(random.choices(string.ascii_letters + string.digits + ' ', k=length))


def make_pdf(filename, password=None):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for _ in range(5):
        pdf.cell(200, 10, txt=random_text(50), ln=True)
    # Get PDF as bytes
    pdf_bytes = pdf.output(dest='S').encode('latin1')
    if password:
        # PyPDF2 can only encrypt an existing PDF
        from PyPDF2 import PdfWriter
        import io
        input_pdf = io.BytesIO(pdf_bytes)
        writer = PdfWriter()
        writer.append(input_pdf)
        writer.encrypt(password)
        with open(filename, "wb") as out_file:
            writer.write(out_file)
    else:
        with open(filename, "wb") as out_file:
            out_file.write(pdf_bytes)

def make_docx(filename, password=None):
    doc = Document()
    doc.add_heading('Random Document', 0)
    for _ in range(5):
        doc.add_paragraph(random_text(100))
    doc.save(filename)
    if password:
        protected_filename = filename.replace(".docx", "_protected.docx")
        with open(filename, 'rb') as f, open(protected_filename, 'wb') as pf:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.encrypt(password, pf)
        os.remove(filename)
        os.rename(protected_filename, filename)

def make_xlsx(filename, password=None):
    output = BytesIO()
    workbook = xlsxwriter.Workbook(output, {'in_memory': True})
    worksheet = workbook.add_worksheet()
    for r in range(10):
        for c in range(5):
            worksheet.write(r, c, random_text(10))
    workbook.close()
    output.seek(0)
    with open(filename, 'wb') as f:
        f.write(output.read())
    if password:
        protected_filename = filename.replace(".xlsx", "_protected.xlsx")
        with open(filename, 'rb') as f, open(protected_filename, 'wb') as pf:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.encrypt(password, pf)
        os.remove(filename)
        os.rename(protected_filename, filename)

def make_pptx(filename, password=None):
    prs = Presentation()
    for _ in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = random_text(10)
        content = random_text(50)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    prs.save(filename)
    if password:
        protected_filename = filename.replace(".pptx", "_protected.pptx")
        with open(filename, 'rb') as f, open(protected_filename, 'wb') as pf:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.encrypt(password, pf)
        os.remove(filename)
        os.rename(protected_filename, filename)


def main():
    password = "Test123!"  # Example password

    files = [
        (make_pdf, ".pdf"),
        (make_docx, ".docx"),
        (make_xlsx, ".xlsx"),
        (make_pptx, ".pptx"),
    ]

    # Not password protected
    for func, suffix in files:
        filename = f'sample_nopass{suffix}'
        func(filename)

    # Password protected
    for func, suffix in files:
        filename = f'sample_pass{suffix}'
        func(filename, password=password)

    print("Files generated in output/ directory.")


if __name__ == '__main__':
    main()
